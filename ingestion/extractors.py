"""Plain-text extraction from many document types.

Each extractor returns a list of ``Block`` sections (heading + optional page +
text). The API's document ingestion chunks these blocks and embeds them.

Targeted, dependency-light parsers are used per format for reliability:
  - PDF            -> PyMuPDF (via ingestion.parsers)
  - DOCX           -> python-docx
  - PPTX           -> python-pptx
  - XLSX           -> openpyxl
  - CSV / HTML / JSON / RTF -> stdlib
  - everything else that decodes as UTF-8 text -> treated as plain text
"""

from __future__ import annotations

import csv as _csv
import io
import json as _json
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import TypedDict


class Block(TypedDict):
    heading: str
    page_number: int | None
    text: str


# Extensions handled as generic UTF-8 text (chunked by Markdown headings).
PLAINTEXT_EXTS = {
    ".txt", ".text", ".md", ".markdown", ".mdown", ".rst", ".log", ".tex",
    ".org", ".adoc", ".asciidoc",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".kt", ".c", ".h", ".cpp",
    ".hpp", ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".scala", ".sh",
    ".ps1", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env", ".sql", ".r",
}

# All extensions the uploader accepts.
SUPPORTED_EXTS = PLAINTEXT_EXTS | {
    ".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".html", ".htm", ".json", ".rtf",
}


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTS


# ---------------------------------------------------------------------------
# Format-specific extractors
# ---------------------------------------------------------------------------

def _pdf_blocks(data: bytes, stem: str) -> list[Block]:
    import tempfile

    from ingestion.parsers import parse_pdf

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
        tf.write(data)
        tmp = Path(tf.name)
    try:
        pages, _ = parse_pdf(tmp)
    finally:
        tmp.unlink(missing_ok=True)
    return [
        {"heading": stem, "page_number": p["page_number"], "text": p["text"]}
        for p in pages
        if p["text"].strip()
    ]


def _docx_blocks(data: bytes, stem: str) -> list[Block]:
    from docx import Document

    doc = Document(io.BytesIO(data))
    blocks: list[Block] = []
    heading = stem
    buf: list[str] = []

    def flush() -> None:
        text = "\n".join(buf).strip()
        if text:
            blocks.append({"heading": heading, "page_number": None, "text": text})

    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading") and para.text.strip():
            flush()
            buf = []
            heading = para.text.strip()
        elif para.text.strip():
            buf.append(para.text)
    flush()

    # Append tables as Markdown at the end (order relative to text is approximate).
    for table in doc.tables:
        rows = [
            "| " + " | ".join(c.text.strip() for c in row.cells) + " |"
            for row in table.rows
        ]
        if rows:
            blocks.append(
                {"heading": stem, "page_number": None, "text": "\n".join(rows)}
            )

    if not blocks:
        joined = "\n".join(p.text for p in doc.paragraphs).strip()
        if joined:
            blocks.append({"heading": stem, "page_number": None, "text": joined})
    return blocks


def _pptx_blocks(data: bytes, stem: str) -> list[Block]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    blocks: list[Block] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        title: str | None = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
        except (AttributeError, ValueError):
            pass
        body = "\n".join(texts).strip()
        if body:
            blocks.append(
                {"heading": title or f"Slide {i}", "page_number": i, "text": body}
            )
    return blocks


def _rows_to_markdown(rows: list[list[str]]) -> str:
    header = rows[0]
    width = len(header)
    md = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        cells = (list(r) + [""] * width)[:width]
        md.append("| " + " | ".join(c.strip() for c in cells) + " |")
    return "\n".join(md)


def _xlsx_blocks(data: bytes, stem: str) -> list[Block]:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[Block] = []
    for ws in wb.worksheets:
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                rows.append(cells)
        if rows:
            blocks.append(
                {"heading": ws.title or stem, "page_number": None,
                 "text": _rows_to_markdown(rows)}
            )
    return blocks


def _csv_blocks(data: bytes, stem: str) -> list[Block]:
    text = data.decode("utf-8-sig", errors="replace")
    rows = [r for r in _csv.reader(io.StringIO(text)) if any(c.strip() for c in r)]
    if not rows:
        return []
    return [{"heading": stem, "page_number": None, "text": _rows_to_markdown(rows)}]


class _HTMLTextExtractor(HTMLParser):
    _BLOCK_TAGS = {"p", "br", "div", "li", "tr", "h1", "h2", "h3", "h4", "h5", "section"}

    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []
        self._skip = 0
        self._in_title = False
        self.title: str | None = None

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style", "noscript"):
            self._skip += 1
        elif tag == "title":
            self._in_title = True
        elif tag in self._BLOCK_TAGS:
            self.parts.append("\n")

    def handle_endtag(self, tag):
        if tag in ("script", "style", "noscript") and self._skip:
            self._skip -= 1
        elif tag == "title":
            self._in_title = False

    def handle_data(self, data):
        if self._skip:
            return
        if self._in_title:
            self.title = (self.title or "") + data
        stripped = data.strip()
        if stripped:
            self.parts.append(stripped + " ")


def _html_blocks(data: bytes, stem: str) -> list[Block]:
    parser = _HTMLTextExtractor()
    parser.feed(data.decode("utf-8", errors="replace"))
    text = "".join(parser.parts)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    heading = (parser.title or stem).strip() or stem
    return [{"heading": heading, "page_number": None, "text": text}]


def _json_blocks(data: bytes, stem: str) -> list[Block]:
    raw = data.decode("utf-8", errors="replace")
    try:
        text = _json.dumps(_json.loads(raw), indent=2, ensure_ascii=False)
    except ValueError:
        text = raw
    return [{"heading": stem, "page_number": None, "text": text}] if text.strip() else []


def _rtf_blocks(data: bytes, stem: str) -> list[Block]:
    raw = data.decode("latin-1", errors="replace")
    raw = re.sub(r"\\'[0-9a-fA-F]{2}", "", raw)          # hex escapes
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)         # control words
    raw = raw.replace("{", "").replace("}", "")
    text = re.sub(r"\n{3,}", "\n\n", raw).strip()
    return [{"heading": stem, "page_number": None, "text": text}] if text else []


def _text_blocks(data: bytes, stem: str) -> list[Block]:
    text = data.decode("utf-8", errors="replace")
    # Split on Markdown headings so each chunk keeps a meaningful heading.
    sections = re.split(r"(?=^#{1,3} )", text, flags=re.MULTILINE)
    sections = [s for s in sections if s.strip()] or [text]
    blocks: list[Block] = []
    for section in sections:
        m = re.match(r"^(#{1,3}) (.+)", section)
        heading = m.group(2).strip() if m else stem
        if section.strip():
            blocks.append({"heading": heading, "page_number": None, "text": section})
    return blocks


_DISPATCH = {
    ".pdf": _pdf_blocks,
    ".docx": _docx_blocks,
    ".pptx": _pptx_blocks,
    ".xlsx": _xlsx_blocks,
    ".csv": _csv_blocks,
    ".html": _html_blocks,
    ".htm": _html_blocks,
    ".json": _json_blocks,
    ".rtf": _rtf_blocks,
}


def extract_blocks(filename: str, data: bytes) -> list[Block]:
    """Extract heading/page/text sections from a document of any supported type.

    Args:
        filename: Original filename (its extension selects the parser).
        data: Raw file bytes.

    Returns:
        A list of Block sections. Unknown extensions fall back to UTF-8 text.
    """
    ext = Path(filename).suffix.lower()
    stem = Path(filename).stem
    handler = _DISPATCH.get(ext)
    if handler is not None:
        return handler(data, stem)
    return _text_blocks(data, stem)
